import os
import tempfile
import numpy as np
import faiss
import requests
import pandas as pd
from sentence_transformers import SentenceTransformer
from langchain.prompts import PromptTemplate
from langchain.chains import RetrievalQA
from langchain_community.vectorstores import FAISS
# from langchain.vectorstores import FAISS
from langchain.embeddings import HuggingFaceEmbeddings
import streamlit as st
import fitz  # PyMuPDF for PDF handling
from docx import Document
from pptx import Presentation
from pyngrok import ngrok
from langchain.schema import Document
from docx import Document as DocxDocument
from langchain.docstore import InMemoryDocstore
from langchain.vectorstores import FAISS as LangchainFAISS

from langchain.docstore.document import Document as LangchainDocument



# Set environment variables
os.environ["REPLICATE_API_TOKEN"] = "r8_CGHRBs7InuuXep74Cg3NMx2kFYYvZJo3mNTSH"
os.environ["HUGGINGFACE_TOKEN"] = "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"

# Initialize models
EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
embeddings_model = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)

# Define utility functions
def extract_text_from_pdf(pdf_file):
    """Extract text from a PDF file."""
    text = ""
    with fitz.open(pdf_file) as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_docx(docx_file):
    """Extract text from a DOCX file."""
    doc = DocxDocument(docx_file)  # Use the aliased import
    return "\n".join([para.text for para in doc.paragraphs])


def extract_text_from_pptx(pptx_file):
    """Extract text from a PPTX file."""
    presentation = Presentation(pptx_file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text

def extract_text_from_xlsx(xlsx_file):
    """Extract text from an XLSX file."""
    df = pd.read_excel(xlsx_file)
    return df.to_string()

def extract_text_from_url(url):
    """Fetch text from a URL."""
    try:
        response = requests.get(url)
        response.raise_for_status()
        return response.text
    except requests.exceptions.RequestException as e:
        st.error(f"Error fetching URL content: {e}")
        return ""

# Define Llama 2 generation function (Simulating the actual model response)
def llama2_generate(prompt, top_p=1, temperature=0.75, max_new_tokens=800):
    """Generate text using Llama 2 (mocked for now)."""
    input = {
        "top_p": top_p,
        "prompt": prompt,
        "temperature": temperature,
        "max_new_tokens": max_new_tokens
    }
    # Mocked response for simplicity
    return f"Generated response based on: {prompt}"

# Main Streamlit application
def main():
    st.title("Llama2-7B Q&A System")

    # Sidebar inputs
    st.sidebar.header("Upload or Provide Input")
    uploaded_files = st.sidebar.file_uploader("Upload files", type=["txt", "pdf", "docx", "pptx", "xlsx"], accept_multiple_files=True)
    web_links = st.sidebar.text_area("Or provide web links (comma-separated):", "").split(',')

    # Allow user to limit number of uploaded files
    num_files_limit = st.sidebar.number_input("Select number of files to upload (Max 10)", min_value=1, max_value=10, value=5)

    chunk_size = st.sidebar.slider("Set Chunk Size", 500, 2000, 1000)

    documents = []

    # Process uploaded files
    if uploaded_files:
        temp_dir = tempfile.TemporaryDirectory()

        for uploaded_file in uploaded_files[:num_files_limit]:  # Limit the number of files uploaded
            file_path = os.path.join(temp_dir.name, uploaded_file.name)
            with open(file_path, "wb") as f:
                f.write(uploaded_file.read())

            if uploaded_file.name.endswith(".txt"):
                with open(file_path, "r", encoding="utf-8") as f:
                    documents.append(f.read())
            elif uploaded_file.name.endswith(".pdf"):
                documents.append(extract_text_from_pdf(file_path))
            elif uploaded_file.name.endswith(".docx"):
                documents.append(extract_text_from_docx(file_path))
            elif uploaded_file.name.endswith(".pptx"):
                documents.append(extract_text_from_pptx(file_path))
            elif uploaded_file.name.endswith(".xlsx"):
                documents.append(extract_text_from_xlsx(file_path))

        temp_dir.cleanup()

    # Process web links
    for web_link in web_links:
        if web_link.strip():  # Skip empty links
            documents.append(extract_text_from_url(web_link.strip()))

    if documents:
        # Split documents into chunks without overlap
        chunks = [doc[i:i+chunk_size] for doc in documents for i in range(0, len(doc), chunk_size)]
        docs = [Document(page_content=chunk, metadata={"source": "uploaded file"}) for chunk in chunks]

        # Generate embeddings
        embeddings = embeddings_model.embed_documents(chunks)

        # Convert embeddings to NumPy array (shape: num_embeddings x embedding_dim)
        embeddings_array = np.array(embeddings)
        dimension = embeddings_array.shape[1]
        faiss_index = faiss.IndexFlatL2(dimension)
        faiss_index.add(embeddings_array.astype(np.float32))

        # Create a docstore
        docstore = InMemoryDocstore({i: docs[i] for i in range(len(docs))})

        # Create index to document mapping
        index_to_docstore_id = {i: str(i) for i in range(len(docs))}

        # Prepare LangChain documents
        vector_store = LangchainFAISS.from_documents(docs, embeddings_model)

        # Set up retriever
        retriever = vector_store.as_retriever()

        # Chat interface
        if "chat_history" not in st.session_state:
            st.session_state.chat_history = []

        st.header("Ask a Question")
        question = st.text_input("Your Question:")
        format_choice = st.selectbox("Answer Format:", ["Default", "Bullet Points", "Summary", "Specific Length"])
        word_limit = st.number_input("Word/Character Limit:", min_value=10, max_value=500, value=100) if format_choice == "Specific Length" else None

        if st.button("Get Answer") and question:
            context_docs = retriever.get_relevant_documents(question)
            context = " ".join([doc.page_content for doc in context_docs])
            prompt = f"Answer the question based on the context: {context}\n\nQuestion: {question}"

            if format_choice == "Bullet Points":
                prompt += "\nPlease provide the answer as bullet points."
            elif format_choice == "Summary":
                prompt += "\nPlease provide a concise summary."
            elif format_choice == "Specific Length":
                prompt += f"\nPlease limit the answer to {word_limit} words."

            answer = llama2_generate(prompt)

            # Save chat history
            st.session_state.chat_history.append({"question": question, "answer": answer})

            st.write("Answer:", answer)

        # Display chat history
        if st.session_state.chat_history:
            st.header("Chat History")
            for entry in st.session_state.chat_history:
                st.write(f"**Q:** {entry['question']}\n**A:** {entry['answer']}\n")

    else:
        st.write("Upload documents or provide a web link to begin.")

if __name__ == "__main__":
    main()